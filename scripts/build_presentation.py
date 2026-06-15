#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Génère une présentation PowerPoint minimaliste et clean pour le projet
« Analyse de radiographies pulmonaires COVID-19 ».

Design épuré, mais contenu COMPLET : toutes les étapes réelles du projet
(exploration, features, prétraitement, baselines ML, transfer learning,
interprétabilité). EfficientNetB0 retenu. Chiffres repris du PPTX actuel.

Usage:
    python3 scripts/build_presentation.py
Sortie:
    Presentation_COVID19_minimaliste.pptx (racine du repo)
"""

import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# --------------------------------------------------------------------------- #
# Chemins
# --------------------------------------------------------------------------- #
ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
FIG = os.path.join(ROOT, "reports", "figures")
REP = os.path.join(ROOT, "reports")
OUT = os.path.join(ROOT, "Presentation_COVID19_minimaliste.pptx")


def fig(name):
    return os.path.join(FIG, name)


def rep(name):
    return os.path.join(REP, name)


# --------------------------------------------------------------------------- #
# Système de design
# --------------------------------------------------------------------------- #
INK = RGBColor(0x16, 0x16, 0x1D)      # texte principal (encre / noir)
ACCENT = RGBColor(0x16, 0x16, 0x1D)   # monochrome : accent = encre (aucune couleur)
MUTED = RGBColor(0x8A, 0x8A, 0x99)    # légendes / kicker (gris)
HAIR = RGBColor(0xE4, 0xE4, 0xE8)     # filets fins
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xC9, 0xC9, 0xD2)    # gris clair

FONT = "Arial"

# Géométrie 16:9
EMU_W = Inches(13.333)
EMU_H = Inches(7.5)
MARGIN = Inches(0.9)
CONTENT_W = EMU_W - 2 * MARGIN

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]

_page = {"n": 0}


# --------------------------------------------------------------------------- #
# Helpers bas niveau
# --------------------------------------------------------------------------- #
def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _set_run(run, text, size, color, bold=False, italic=False, spacing=None):
    run.text = text
    f = run.font
    f.name = FONT
    f.size = Pt(size)
    f.color.rgb = color
    f.bold = bold
    f.italic = italic
    if spacing is not None:
        rPr = run._r.get_or_add_rPr()
        rPr.set("spc", str(int(spacing * 100)))


def textbox(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tb, tf


def _add_break(p):
    br = p._p.makeelement(qn("a:br"), {})
    p._p.append(br)


def add_para(tf, text, size, color, bold=False, italic=False, spacing=None,
             align=PP_ALIGN.LEFT, space_after=6, space_before=0, line=None, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    if line is not None:
        p.line_spacing = line
    parts = text.split("\n")
    for i, part in enumerate(parts):
        if i > 0:
            _add_break(p)
        r = p.add_run()
        _set_run(r, part, size, color, bold=bold, italic=italic, spacing=spacing)
    return p


def line(slide, left, top, width, color=HAIR, weight=1.2):
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Emu(int(Pt(weight))))
    ln.fill.solid()
    ln.fill.fore_color.rgb = color
    ln.line.fill.background()
    ln.shadow.inherit = False
    return ln


def rect(slide, left, top, width, height, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def footer(slide, dark=False):
    _page["n"] += 1
    col = LIGHT if dark else MUTED
    tb, tf = textbox(slide, MARGIN, EMU_H - Inches(0.55),
                     CONTENT_W, Inches(0.3), anchor=MSO_ANCHOR.MIDDLE)
    add_para(tf, "Radiographies pulmonaires COVID-19", 9, col, first=True)
    nb, nf = textbox(slide, EMU_W - MARGIN - Inches(1.0), EMU_H - Inches(0.55),
                     Inches(1.0), Inches(0.3), anchor=MSO_ANCHOR.MIDDLE)
    add_para(nf, f"{_page['n']:02d}", 9, col, align=PP_ALIGN.RIGHT, first=True)


def kicker_title(slide, kicker, title, dark=False):
    top = MARGIN
    tb, tf = textbox(slide, MARGIN, top, CONTENT_W, Inches(0.3))
    add_para(tf, kicker.upper(), 11, MUTED, bold=True, spacing=2.2, first=True)
    tt, ttf = textbox(slide, MARGIN, top + Inches(0.34), CONTENT_W, Inches(1.0))
    add_para(ttf, title, 31, WHITE if dark else INK, bold=False, line=1.02, first=True)
    line(slide, MARGIN, top + Inches(1.28), Inches(1.1), color=ACCENT, weight=2.4)


def slide_blank(dark=False):
    s = prs.slides.add_slide(BLANK)
    _bg(s, INK if dark else WHITE)
    return s


def fit_image(slide, path, box_left, box_top, box_w, box_h, caption=None):
    with Image.open(path) as im:
        iw, ih = im.size
    ratio = iw / ih
    box_ratio = box_w / box_h
    if ratio > box_ratio:
        w = box_w
        h = int(box_w / ratio)
    else:
        h = box_h
        w = int(box_h * ratio)
    left = box_left + (box_w - w) // 2
    top = box_top + (box_h - h) // 2
    slide.shapes.add_picture(path, left, top, width=w, height=h)
    if caption:
        cb, cf = textbox(slide, box_left, box_top + box_h + Inches(0.05),
                         box_w, Inches(0.3))
        add_para(cf, caption, 10, MUTED, italic=True, align=PP_ALIGN.CENTER, first=True)


# --------------------------------------------------------------------------- #
# Gabarits de slides
# --------------------------------------------------------------------------- #
def title_slide():
    s = slide_blank()
    tb, tf = textbox(s, MARGIN, Inches(2.0), Inches(10.8), Inches(2.6))
    add_para(tf, "PROJET DATA SCIENCE · CLASSIFICATION D'IMAGES MÉDICALES",
             11, MUTED, bold=True, spacing=2.0, first=True, space_after=14)
    add_para(tf, "Analyse de radiographies", 46, INK, line=1.0, space_after=0)
    add_para(tf, "pulmonaires COVID-19", 46, INK, line=1.0, space_after=0)
    line(s, MARGIN, Inches(4.7), Inches(1.4), color=INK, weight=2.6)
    ab, af = textbox(s, MARGIN, Inches(5.2), Inches(11.0), Inches(1.4))
    add_para(af, "Souad Sadoun   ·   Théo Daguier   ·   Mohamed Amghar",
             16, INK, bold=True, first=True, space_after=6)
    add_para(af, "Formation DataScientest   ·   Mentor : Nicolas Mormiche   ·   Juin 2026",
             13, MUTED)
    return s


def section_divider(kicker, title, step=None):
    s = slide_blank()
    tb, tf = textbox(s, MARGIN, Inches(2.7), CONTENT_W, Inches(2.0))
    add_para(tf, kicker.upper(), 12, MUTED, bold=True, spacing=2.6, first=True, space_after=16)
    add_para(tf, title, 42, INK, line=1.02)
    line(s, MARGIN, Inches(4.55), Inches(1.4), color=INK, weight=2.6)
    if step:
        nb, nf = textbox(s, EMU_W - MARGIN - Inches(3.0), Inches(2.75),
                         Inches(3.0), Inches(0.5))
        add_para(nf, step, 13, MUTED, align=PP_ALIGN.RIGHT, first=True)
    footer(s)
    return s


def content_slide(kicker, title, bullets, lead=None):
    s = slide_blank()
    kicker_title(s, kicker, title)
    top = Inches(2.4)
    if lead:
        lb, lf = textbox(s, MARGIN, top, CONTENT_W, Inches(0.9))
        add_para(lf, lead, 18, INK, line=1.25, first=True)
        top = Inches(3.35)
    bb, bf = textbox(s, MARGIN, top, CONTENT_W, Inches(3.6))
    for i, b in enumerate(bullets):
        add_para(bf, b, 17, INK, line=1.2, space_after=12, first=(i == 0))
    footer(s)
    return s


def content_with_image(kicker, title, bullets, image, caption=None, lead=None):
    s = slide_blank()
    kicker_title(s, kicker, title)
    col_w = Inches(5.1)
    top = Inches(2.4)
    if lead:
        lb, lf = textbox(s, MARGIN, top, col_w, Inches(0.9))
        add_para(lf, lead, 16, INK, line=1.25, first=True)
        top = top + Inches(1.0)
    bb, bf = textbox(s, MARGIN, top, col_w, Inches(3.6))
    for i, b in enumerate(bullets):
        add_para(bf, b, 15, INK, line=1.2, space_after=10, first=(i == 0))
    img_left = MARGIN + col_w + Inches(0.4)
    img_w = EMU_W - MARGIN - img_left
    fit_image(s, image, img_left, Inches(2.3), img_w, Inches(4.1), caption=caption)
    footer(s)
    return s


def metric_slide(kicker, title, metrics, note=None, sub=None):
    s = slide_blank()
    kicker_title(s, kicker, title)
    top = Inches(2.8)
    n = len(metrics)
    gap = Inches(0.5)
    total_gap = gap * (n - 1)
    card_w = (CONTENT_W - total_gap) / n
    for i, (val, lab) in enumerate(metrics):
        left = MARGIN + i * (card_w + gap)
        accent = (i == 0)
        vb, vf = textbox(s, left, top, card_w, Inches(1.6), anchor=MSO_ANCHOR.BOTTOM)
        add_para(vf, val, 68 if n <= 3 else 52, ACCENT if accent else INK,
                 bold=True, first=True)
        lb, lf = textbox(s, left, top + Inches(1.72), card_w, Inches(0.9))
        add_para(lf, lab, 14, MUTED, line=1.15, first=True)
        line(s, left, top + Inches(1.64), Inches(0.7),
             color=ACCENT if accent else HAIR, weight=2.2)
    if sub:
        sb, sf = textbox(s, MARGIN, Inches(5.3), CONTENT_W, Inches(0.9))
        add_para(sf, sub, 16, INK, line=1.25, first=True)
    if note:
        nb, nf = textbox(s, MARGIN, Inches(6.15), CONTENT_W, Inches(0.7))
        add_para(nf, note, 13, MUTED, italic=True, first=True)
    footer(s)
    return s


def image_slide(kicker, title, image, caption=None, note=None):
    s = slide_blank()
    kicker_title(s, kicker, title)
    box_top = Inches(2.35)
    box_h = Inches(4.0) if note else Inches(4.25)
    fit_image(s, image, MARGIN, box_top, CONTENT_W, box_h, caption=caption)
    if note:
        nb, nf = textbox(s, MARGIN, Inches(6.7), CONTENT_W, Inches(0.5))
        add_para(nf, note, 13, MUTED, italic=True, align=PP_ALIGN.CENTER, first=True)
    footer(s)
    return s


def two_image_slide(kicker, title, img1, cap1, img2, cap2, note=None):
    s = slide_blank()
    kicker_title(s, kicker, title)
    half = (CONTENT_W - Inches(0.5)) / 2
    fit_image(s, img1, MARGIN, Inches(2.35), half, Inches(3.9), caption=cap1)
    fit_image(s, img2, MARGIN + half + Inches(0.5), Inches(2.35), half, Inches(3.9), caption=cap2)
    if note:
        nb, nf = textbox(s, MARGIN, Inches(6.75), CONTENT_W, Inches(0.5))
        add_para(nf, note, 13, MUTED, italic=True, align=PP_ALIGN.CENTER, first=True)
    footer(s)
    return s


def four_image_slide(kicker, title, items, note=None):
    s = slide_blank()
    kicker_title(s, kicker, title)
    gap = Inches(0.4)
    cell_w = (CONTENT_W - gap) / 2
    cell_h = Inches(1.9)
    top0 = Inches(2.3)
    for i, (img, cap) in enumerate(items[:4]):
        r, c = divmod(i, 2)
        left = MARGIN + c * (cell_w + gap)
        top = top0 + r * (cell_h + Inches(0.45))
        fit_image(s, img, left, top, cell_w, cell_h, caption=cap)
    if note:
        nb, nf = textbox(s, MARGIN, EMU_H - Inches(0.9), CONTENT_W, Inches(0.4))
        add_para(nf, note, 12, MUTED, italic=True, align=PP_ALIGN.CENTER, first=True)
    footer(s)
    return s


def table_slide(kicker, title, headers, rows, highlight=None, note=None, top=Inches(2.5)):
    s = slide_blank()
    kicker_title(s, kicker, title)
    n_cols = len(headers)
    n_rows = len(rows) + 1
    tbl_w = CONTENT_W
    row_h = Inches(0.58)
    tbl_h = row_h * n_rows
    gtbl = s.shapes.add_table(n_rows, n_cols, MARGIN, top, tbl_w, tbl_h).table
    first_w = int(tbl_w * 0.30)
    other_w = int((tbl_w - first_w) / (n_cols - 1))
    gtbl.columns[0].width = first_w
    for c in range(1, n_cols):
        gtbl.columns[c].width = other_w
    tblPr = gtbl._tbl.tblPr
    tblPr.set("firstRow", "0")
    tblPr.set("bandRow", "0")

    def style_cell(cell, text, size, color, bold=False, fill=None, align=PP_ALIGN.LEFT):
        cell.margin_left = Inches(0.12)
        cell.margin_right = Inches(0.08)
        cell.margin_top = Inches(0.03)
        cell.margin_bottom = Inches(0.03)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.fill.solid()
        cell.fill.fore_color.rgb = WHITE if fill is None else fill
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        _set_run(r, text, size, color, bold=bold)

    for c, h in enumerate(headers):
        style_cell(gtbl.cell(0, c), h, 13, MUTED, bold=True,
                   align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER)
    for ri, row in enumerate(rows):
        is_hl = (highlight is not None and ri == highlight)
        fill = ACCENT if is_hl else WHITE
        txt_col = WHITE if is_hl else INK
        for c, val in enumerate(row):
            style_cell(gtbl.cell(ri + 1, c), val, 14, txt_col,
                       bold=(c == 0 or is_hl), fill=fill,
                       align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER)
    if note:
        nb, nf = textbox(s, MARGIN, top + tbl_h + Inches(0.3), CONTENT_W, Inches(1.6))
        add_para(nf, note, 14, INK, line=1.3, first=True)
    footer(s)
    return s


# Rapport de classification par classe (gabarit dérivé de table_slide)
PCLASS_HEADERS = ["Classe", "Précision", "Rappel", "F1-score", "Support"]


def per_class_slide(kicker, title, rows, note=None):
    return table_slide(kicker, title, PCLASS_HEADERS, rows, note=note, top=Inches(2.5))


# --------------------------------------------------------------------------- #
# Construction du deck
# --------------------------------------------------------------------------- #
def build():
    # ---------------- Intro ----------------
    title_slide()

    content_slide(
        "Contexte", "Pourquoi ce projet ?",
        bullets=[
            "—  La pandémie de COVID-19 a mis sous tension les systèmes de santé.",
            "—  La radiographie thoracique est un examen rapide, peu coûteux et largement disponible.",
            "—  Une IA peut assister le radiologue dans le tri et la priorisation des cas, sans s'y substituer.",
        ],
        lead="La radiographie pulmonaire, lue par un modèle d'apprentissage, devient un outil d'aide à la décision.",
    )

    content_slide(
        "Objectif", "Classer automatiquement 4 pathologies",
        bullets=[
            "—  COVID-19",
            "—  Normal (poumon sain)",
            "—  Lung Opacity (opacité pulmonaire)",
            "—  Viral Pneumonia (pneumonie virale)",
        ],
        lead="Objectif : un modèle de Deep Learning qui classe une radiographie thoracique en l'une de 4 catégories.",
    )

    content_slide(
        "Démarche", "Notre fil conducteur en 4 étapes",
        bullets=[
            "1.  Comprendre & explorer les données — distribution, qualité, caractéristiques d'images.",
            "2.  Prétraiter & enrichir — normalisation, débruitage, extraction de features, augmentation.",
            "3.  Modéliser — baselines Machine Learning, puis Transfer Learning fine-tuné.",
            "4.  Interpréter & décider — Grad-CAM, SHAP, choix du modèle final.",
        ],
    )

    # ---------------- Étape 1 : Exploration ----------------
    section_divider("Étape 1", "Compréhension & exploration des données")

    metric_slide(
        "Données", "COVID-19 Radiography Database",
        metrics=[("21 165", "images PNG\nau total"),
                 ("4", "classes\nde pathologies"),
                 ("7,6×", "écart majoritaire /\nminoritaire")],
        sub="Normal 10 192  ·  Lung Opacity 6 012  ·  COVID 3 616  ·  Viral Pneumonia 1 345",
        note="Source : Kaggle. Images PNG en niveaux de gris.",
    )

    image_slide(
        "Exploration", "Une répartition fortement déséquilibrée",
        rep("Pourcentage d'images par catégorie.png"),
        caption="Part de chaque classe dans le jeu de données",
        note="Le déséquilibre marqué (Normal ≈ 48 %, Viral Pneumonia ≈ 6 %) conditionne toute la méthodologie.",
    )

    content_slide(
        "Exploration", "Un jeu de données propre",
        bullets=[
            "—  Recherche de doublons  →  aucun doublon détecté.",
            "—  Contrôle des dimensions et du format  →  PNG uniforme, images homogènes.",
            "—  Vérification de la qualité des images  →  toutes exploitables.",
            "—  Analyse de la répartition des classes  →  déséquilibre identifié et traité.",
        ],
        lead="Vérifications réalisées avant toute modélisation :",
    )

    image_slide(
        "Exploration", "Un exemple par classe",
        fig("exemple_par_classe.png"),
        caption="COVID · Lung Opacity · Normal · Viral Pneumonia",
    )

    content_with_image(
        "Exploration", "Caractéristiques des images",
        bullets=[
            "—  Intensité moyenne",
            "—  Écart-type",
            "—  Contraste",
            "—  Entropie",
            "—  Gradient des contours",
        ],
        image=rep("hist_moyennes.png"),
        caption="Distribution de l'intensité moyenne par classe",
        lead="Plusieurs métriques quantitatives ont été extraites de chaque image :",
    )

    content_with_image(
        "Exploration", "Corrélation des caractéristiques",
        bullets=[
            "—  Intensité, contraste et écart-type sont fortement corrélés.",
            "—  Entropie et gradient apportent une information complémentaire.",
            "—  Conclusion : on peut réduire la dimension (ACP) sans perdre l'essentiel.",
        ],
        image=rep("Comparaison_preprocessing_features.png"),
        caption="Comparaison des features de prétraitement",
    )

    image_slide(
        "Exploration", "Visualisation par ACP",
        fig("pca_corrigee.png"),
        caption="Projection en composantes principales",
        note="Séparation partielle des classes et identification de quelques images atypiques (outliers).",
    )

    # ---------------- Étape 2 : Prétraitement ----------------
    section_divider("Étape 2", "Prétraitement & ingénierie des features")

    content_slide(
        "Prétraitement", "Un pipeline en 4 temps",
        bullets=[
            "1.  Normalisation  —  niveaux de gris, redimensionnement 224×224, pixels dans [0, 1].",
            "2.  Débruitage  —  filtre gaussien, filtres morphologiques.",
            "3.  Extraction de features  —  Canny, Sobel, Laplacien, érosion / dilatation.",
            "4.  Augmentation  —  rotation, zoom, translation, flip horizontal.",
        ],
        lead="Objectif : homogénéiser les entrées et faciliter l'apprentissage des modèles.",
    )

    content_with_image(
        "Prétraitement", "Normalisation des images",
        bullets=[
            "—  Conversion en niveaux de gris.",
            "—  Redimensionnement à 224 × 224 pixels.",
            "—  Normalisation des pixels entre 0 et 1.",
        ],
        image=rep("image_pretraitee_0.png"),
        caption="Image après normalisation",
        lead="Homogénéiser taille et intensité de toutes les images.",
    )

    content_with_image(
        "Prétraitement", "Filtre gaussien (débruitage)",
        bullets=[
            "—  Lisse les images et réduit le bruit de certaines radiographies.",
            "—  Améliore la qualité visuelle.",
            "—  Facilite l'extraction des caractéristiques pulmonaires.",
            "—  Stabilise l'apprentissage des modèles.",
        ],
        image=rep("image_gaussian_blur_0.png"),
        caption="Application d'un Gaussian Blur",
    )

    four_image_slide(
        "Prétraitement", "Extraction de caractéristiques visuelles",
        items=[
            (rep("image_canny_0.png"), "Canny — contours"),
            (rep("image_sobel_0.png"), "Sobel — gradients"),
            (rep("image_laplacian_0.png"), "Laplacien — détails fins"),
            (rep("image_erosion_0.png"), "Érosion — morphologie"),
        ],
        note="Ces filtres explorent la structure des images médicales et mettent en évidence différentes caractéristiques.",
    )

    content_slide(
        "Prétraitement", "Data augmentation",
        bullets=[
            "—  Rotation des images (±10 %).",
            "—  Zoom (±10 %) et translation (±5 %).",
            "—  Flip horizontal uniquement — pas de flip vertical (l'anatomie gauche/droite est informative).",
            "—  Effet : plus de données artificielles, meilleure généralisation, moins de sur-apprentissage.",
        ],
        lead="Pour améliorer la robustesse du modèle au cadrage et aux variations.",
    )

    table_slide(
        "Prétraitement", "Gestion du déséquilibre — poids de classe",
        headers=["Classe", "Effectif", "Poids appliqué"],
        rows=[
            ["COVID-19", "3 616", "1,46"],
            ["Lung Opacity", "6 012", "0,88"],
            ["Normal", "10 192", "0,52"],
            ["Viral Pneumonia", "1 345", "3,93"],
        ],
        note="Des poids de classe « balanced » (calculés sur le train) sont passés à l'entraînement : "
             "ils donnent plus d'importance aux classes rares, réduisent le biais vers Normal et "
             "améliorent le rappel des pathologies minoritaires.",
    )

    # ---------------- Étape 3 : Baseline ML ----------------
    section_divider("Étape 3", "Baseline Machine Learning")

    metric_slide(
        "Modèle classique", "Random Forest (features HOG)",
        metrics=[("81,95 %", "Accuracy"),
                 ("81,40 %", "Macro F1-score"),
                 ("85,76 %", "Précision macro")],
        note="Meilleur des modèles classiques.",
    )
    per_class_slide(
        "Modèle classique · Random Forest", "Rapport de classification",
        rows=[
            ["COVID", "0,90", "0,70", "0,78", "723"],
            ["Lung Opacity", "0,80", "0,73", "0,76", "1 203"],
            ["Normal", "0,80", "0,92", "0,86", "2 038"],
            ["Viral Pneumonia", "0,94", "0,78", "0,85", "269"],
        ],
        note="Le passage des images en vecteur 1D (flattening) fait perdre l'information spatiale des poumons — "
             "limite intrinsèque des modèles classiques.",
    )

    metric_slide(
        "Modèle classique", "K-Nearest Neighbors (KNN)",
        metrics=[("77,39 %", "Accuracy"),
                 ("77,04 %", "Macro F1-score"),
                 ("77,08 %", "Rappel macro")],
        note="Sensible au bruit et peu adapté aux images en haute dimension.",
    )
    per_class_slide(
        "Modèle classique · KNN", "Rapport de classification",
        rows=[
            ["COVID", "0,69", "0,67", "0,68", "723"],
            ["Lung Opacity", "0,72", "0,69", "0,71", "1 203"],
            ["Normal", "0,82", "0,85", "0,83", "2 038"],
            ["Viral Pneumonia", "0,85", "0,88", "0,86", "269"],
        ],
    )

    metric_slide(
        "Deep Learning", "CNN « maison » (baseline)",
        metrics=[("74,99 %", "Accuracy"),
                 ("72,43 %", "Macro F1-score"),
                 ("72,89 %", "Rappel macro")],
        note="Rappel COVID faible (0,42) : le réseau entraîné de zéro sur-apprend et généralise mal.",
    )
    image_slide(
        "Deep Learning · CNN", "Courbes d'apprentissage — sur-apprentissage",
        rep("loss_accuracy_curves.png"),
        caption="Accuracy & loss, train vs validation",
        note="L'écart train / validation révèle un sur-apprentissage marqué du CNN entraîné de zéro.",
    )

    metric_slide(
        "Deep Learning", "Hybride CNN + Gradient Boosting",
        metrics=[("83,3 %", "Accuracy"),
                 ("84,2 %", "Macro F1-score")],
        sub="Le CNN sert d'extracteur de features, le Gradient Boosting décide.",
        note="Meilleure solution « classique » du projet — avant de passer au Transfer Learning.",
    )
    image_slide(
        "Deep Learning · Hybride", "Matrice de confusion",
        rep("confusion_boosting.png"),
        caption="Hybride CNN + Gradient Boosting",
    )

    # ---------------- Étape 4 : Transfer Learning ----------------
    section_divider("Étape 4", "Transfer Learning")

    content_slide(
        "Transfer Learning", "Réutiliser les connaissances d'ImageNet",
        bullets=[
            "—  Peu de données médicales annotées → on part de backbones pré-entraînés sur ImageNet.",
            "—  Deux architectures de familles différentes : ResNet50 (connexions résiduelles) et EfficientNetB0 (compound scaling).",
            "—  Chaque backbone applique son preprocessing dédié (jamais un /255 artisanal).",
        ],
        lead="Moins de données à apprendre, des features visuelles déjà solides : le bon levier ici.",
    )

    content_slide(
        "Transfer Learning", "Protocole de fine-tuning en 2 phases",
        bullets=[
            "Phase 1 — extraction de features : backbone gelé, on entraîne uniquement la tête de classification (lr = 1e-3).",
            "Phase 2 — fine-tuning : on dégèle le dernier bloc à très faible learning rate (lr = 1e-5).",
            "Point critique : les couches BatchNorm restent gelées — les dégeler ferait diverger l'entraînement.",
            "Garde-fous : poids de classe, EarlyStopping, ModelCheckpoint, ReduceLROnPlateau.",
        ],
    )

    # ResNet50
    metric_slide(
        "Modèle", "ResNet50 — fine-tuné",
        metrics=[("93,32 %", "Accuracy"),
                 ("93,58 %", "Macro F1-score"),
                 ("93,15 %", "Rappel macro")],
        note="23,8 M de paramètres — le poids lourd, très précis.",
    )
    per_class_slide(
        "Modèle · ResNet50", "Rapport de classification",
        rows=[
            ["COVID", "0,96", "0,96", "0,96", "542"],
            ["Lung Opacity", "0,94", "0,87", "0,91", "902"],
            ["Normal", "0,92", "0,96", "0,94", "1 529"],
            ["Viral Pneumonia", "0,94", "0,94", "0,94", "202"],
        ],
    )
    two_image_slide(
        "Modèle · ResNet50", "Confusion & courbes d'apprentissage",
        fig("confusion_resnet50.png"), "Matrice de confusion",
        fig("learning_curves_resnet50.png"), "Courbes train / validation",
        note="Confusions résiduelles surtout entre Lung Opacity et Normal ; sur-apprentissage contenu.",
    )
    image_slide(
        "Modèle · ResNet50", "Courbes ROC & Précision-Rappel",
        fig("roc_pr_resnet50.png"),
        caption="ROC et PR (One-vs-Rest) — ROC-AUC ≈ 0,985",
    )

    # EfficientNetB0
    metric_slide(
        "Modèle", "EfficientNetB0 — fine-tuné",
        metrics=[("89,29 %", "Accuracy"),
                 ("89,61 %", "Macro F1-score"),
                 ("91,34 %", "Rappel macro")],
        note="4,2 M de paramètres seulement — ~5× plus léger que ResNet50 pour une performance proche.",
    )
    per_class_slide(
        "Modèle · EfficientNetB0", "Rapport de classification",
        rows=[
            ["COVID", "0,84", "0,92", "0,88", "542"],
            ["Lung Opacity", "0,86", "0,90", "0,88", "902"],
            ["Normal", "0,94", "0,87", "0,90", "1 529"],
            ["Viral Pneumonia", "0,89", "0,96", "0,92", "202"],
        ],
        note="Excellent rappel sur les classes minoritaires (COVID 0,92 · Viral Pneumonia 0,96).",
    )
    two_image_slide(
        "Modèle · EfficientNetB0", "Confusion & courbes d'apprentissage",
        fig("confusion_efficientnetb0.png"), "Matrice de confusion",
        fig("learning_curves_efficientnetb0.png"), "Courbes train / validation",
    )
    image_slide(
        "Modèle · EfficientNetB0", "Courbes ROC & Précision-Rappel",
        fig("roc_pr_efficientnetb0.png"),
        caption="ROC et PR (One-vs-Rest) — ROC-AUC ≈ 0,984",
    )

    # Comparaison
    table_slide(
        "Comparaison", "Tous les modèles, côte à côte",
        headers=["Modèle", "Accuracy", "Macro F1", "Rappel macro", "Paramètres"],
        rows=[
            ["CNN baseline", "74,99 %", "72,43 %", "72,89 %", "0,1 M"],
            ["Hybride CNN + GB", "83,3 %", "84,2 %", "—", "—"],
            ["EfficientNetB0", "89,29 %", "89,61 %", "91,34 %", "4,2 M"],
            ["ResNet50", "93,32 %", "93,58 %", "93,15 %", "23,8 M"],
        ],
        note="Les deux modèles de Transfer Learning dominent largement les approches classiques.",
    )

    # ---------------- Étape 5 : Interprétabilité ----------------
    section_divider("Étape 5", "Interprétabilité")

    two_image_slide(
        "Interprétabilité", "Grad-CAM : où le modèle regarde-t-il ?",
        fig("gradcam_efficientnetb0_COVID.png"), "EfficientNetB0 — cas COVID",
        fig("gradcam_resnet50_COVID.png"), "ResNet50 — cas COVID",
        note="Rouge = forte contribution. Les activations se concentrent sur les régions pulmonaires pertinentes.",
    )
    four_image_slide(
        "Interprétabilité · Grad-CAM", "EfficientNetB0 — une carte par classe",
        items=[
            (fig("gradcam_efficientnetb0_COVID.png"), "COVID"),
            (fig("gradcam_efficientnetb0_Lung_Opacity.png"), "Lung Opacity"),
            (fig("gradcam_efficientnetb0_Normal.png"), "Normal"),
            (fig("gradcam_efficientnetb0_Viral_Pneumonia.png"), "Viral Pneumonia"),
        ],
    )
    two_image_slide(
        "Interprétabilité · Grad-CAM", "Analyse des erreurs (cas mal classés)",
        fig("gradcam_efficientnetb0_miscls_0.png"), "EfficientNetB0 — erreur",
        fig("gradcam_resnet50_miscls_0.png"), "ResNet50 — erreur",
        note="Sur les cas mal classés, l'attention dérive parfois hors des zones pathologiques — piste d'amélioration.",
    )
    image_slide(
        "Interprétabilité", "SHAP : quantifier l'influence de chaque région",
        fig("shap_efficientnetb0.png"),
        caption="Rouge = augmente la probabilité de la classe · Bleu = la diminue",
    )
    content_slide(
        "Interprétabilité", "Des décisions médicalement cohérentes",
        bullets=[
            "—  Cohérence médicale : les zones importantes correspondent aux régions pulmonaires pertinentes.",
            "—  Décisions explicables : chaque prédiction se justifie par des éléments visuels observables.",
            "—  Complémentarité : Grad-CAM localise, SHAP quantifie l'influence positive ou négative.",
            "—  Confiance : le modèle s'appuie sur la pathologie, pas sur des artefacts d'image.",
        ],
    )

    # ---------------- Conclusion ----------------
    section_divider("Bilan", "Conclusion & perspectives")

    table_slide(
        "Décision finale", "EfficientNetB0, le meilleur compromis",
        headers=["Modèle", "Macro F1", "Paramètres", "Décision"],
        rows=[
            ["CNN baseline", "72,43 %", "0,1 M", "Écarté"],
            ["ResNet50", "93,58 %", "23,8 M", "Alternative très précise"],
            ["EfficientNetB0", "89,61 %", "4,2 M", "Retenu"],
        ],
        highlight=2,
        note="EfficientNetB0 atteint une performance proche de ResNet50 avec ~5× moins de paramètres et "
             "un excellent rappel sur les classes minoritaires : le meilleur compromis performance / coût "
             "pour un outil d'aide au diagnostic déployable. ResNet50 reste l'option si la précision brute prime.",
    )

    content_slide(
        "Bilan", "Difficultés rencontrées",
        bullets=[
            "—  Déséquilibre des classes (Normal 48 % vs Viral 6 %) → poids de classe + augmentation + split stratifié.",
            "—  Perte d'information spatiale en ML (flattening 1D) → résolue par les CNN / Transfer Learning.",
            "—  Sur-apprentissage du CNN seul → approche hybride puis Transfer Learning.",
            "—  Coût de calcul sur 21 165 images → EarlyStopping, ModelCheckpoint, précision mixte.",
        ],
    )

    content_slide(
        "Conclusion", "Objectifs atteints & perspectives",
        bullets=[
            "✓  Transfer Learning et fine-tuning de 2 architectures pré-entraînées.",
            "✓  Interprétabilité approfondie via Grad-CAM et SHAP.",
            "✓  Comparaison rigoureuse de tous les modèles.",
            "→  Perspectives : dataset multi-centrique, architectures récentes (ViT, EfficientNetV2), "
            "données cliniques complémentaires, outil d'aide au diagnostic en temps réel.",
        ],
        lead="Un pipeline complet, du prétraitement à un modèle interprétable et déployable.",
    )

    # Merci
    s = slide_blank()
    tb, tf = textbox(s, MARGIN, Inches(2.9), CONTENT_W, Inches(1.6))
    add_para(tf, "Merci pour votre attention", 40, INK, first=True, space_after=14)
    add_para(tf, "Souad Sadoun  ·  Théo Daguier  ·  Mohamed Amghar", 16, MUTED)
    line(s, MARGIN, Inches(4.5), Inches(1.4), color=INK, weight=2.6)

    prs.save(OUT)
    print(f"OK — {len(prs.slides._sldIdLst)} slides")
    print(f"Écrit : {OUT}")


if __name__ == "__main__":
    build()
